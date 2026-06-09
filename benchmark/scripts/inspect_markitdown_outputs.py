#!/usr/bin/env python3
from __future__ import annotations

import argparse
import csv
import datetime as dt
import hashlib
import multiprocessing
import re
import tarfile
import tomllib
import unicodedata
import zipfile
from pathlib import Path
from queue import Empty
from xml.etree import ElementTree


TEXT_EXTENSIONS = {
    ".csv",
    ".gdoc",
    ".gsheet",
    ".gslides",
    ".htm",
    ".html",
    ".json",
    ".md",
    ".txt",
    ".velja-rules",
    ".yaml",
    ".yml",
}


def main() -> int:
    parser = argparse.ArgumentParser()
    parser.add_argument("--config")
    parser.add_argument("--inventory")
    parser.add_argument("--markdown-root")
    parser.add_argument("--out")
    parser.add_argument("--report")
    parser.add_argument("--source-timeout-seconds", type=int)
    args = parser.parse_args()

    config = load_config(args.config)
    inventory = Path(
        value_from_args_or_config(
            args.inventory,
            config,
            "inspection_inventory",
            config.get("inventory_out", "benchmark/reports/markitdown-inventory.tsv"),
        )
    )
    markdown_root = Path(
        value_from_args_or_config(
            args.markdown_root,
            config,
            "inspection_markdown_root",
            config.get("markdown_output_root", "benchmark/outputs/markitdown"),
        )
    )
    out = Path(
        value_from_args_or_config(
            args.out,
            config,
            "inspection_out",
            "benchmark/reports/markitdown-inspection.tsv",
        )
    )
    report = Path(
        value_from_args_or_config(
            args.report,
            config,
            "inspection_report",
            "benchmark/reports/markitdown-report.md",
        )
    )
    source_timeout_seconds = args.source_timeout_seconds
    if source_timeout_seconds is None:
        source_timeout_seconds = int(config.get("inspection_source_timeout_seconds", "120"))

    rows = read_inventory(inventory)
    ok_rows = [row for row in rows if row["status"] == "ok"]
    inspected = []
    for index, row in enumerate(ok_rows, start=1):
        inspected.append(inspect_row(row, markdown_root, source_timeout_seconds))
        print(f"{index}/{len(ok_rows)} {row['path']}", flush=True)
    out.parent.mkdir(parents=True, exist_ok=True)
    write_tsv(out, inspected)
    append_report(report, inventory, markdown_root, out, rows, inspected)
    print(out)
    print(report)
    return 0


def load_config(config_arg: str | None) -> dict[str, str]:
    config_path = Path(config_arg) if config_arg else Path("benchmark/benchmark.config.toml")
    if not config_path.exists():
        if config_arg:
            raise FileNotFoundError(f"config file not found: {config_path}")
        return {}
    with config_path.open("rb") as file:
        data = tomllib.load(file)
    return {str(key): str(value) for key, value in data.items()}


def value_from_args_or_config(
    arg_value: str | None, config: dict[str, str], key: str, default: str | None = None
) -> str:
    value = arg_value or config.get(key) or default
    if value is None:
        raise ValueError(f"{key} is required")
    return value


def read_inventory(inventory: Path) -> list[dict[str, str]]:
    with inventory.open(encoding="utf-8", newline="") as file:
        return list(csv.DictReader(file, delimiter="\t"))


def inspect_row(
    inventory_row: dict[str, str],
    markdown_root: Path,
    source_timeout_seconds: int,
) -> dict[str, str]:
    source_path = Path(inventory_row["path"])
    markdown_path = markdown_root / f"{safe_name(str(source_path))}.md"
    return inspect_pair_with_timeout(source_path, markdown_path, source_timeout_seconds)


def inspect_pair(source_path: Path, markdown_path: Path) -> dict[str, str]:
    markdown = read_text(markdown_path) if markdown_path.exists() else ""
    markdown_chars = len(markdown)
    markdown_language = classify_language(markdown)

    try:
        source_kind, source_text = extract_source_text(source_path)
    except Exception as error:
        return result_row(
            source_path=source_path,
            markdown_path=markdown_path,
            source_kind=source_path.suffix.lower().lstrip(".") or "unknown",
            source_language="unknown",
            markdown_language=markdown_language,
            language_status="unknown",
            source_chars=0,
            markdown_chars=markdown_chars,
            char_ratio="",
            overlap="",
            content_status="unknown",
            notes=f"{type(error).__name__}: {error}",
        )

    source_chars = len(source_text)
    source_language = classify_language(source_text)
    language_status = compare_language(source_language, markdown_language)
    char_ratio_value = markdown_chars / source_chars if source_chars else 0.0
    overlap_value = overlap_score(source_text, markdown)
    content_status = compare_content(source_chars, markdown_chars, char_ratio_value, overlap_value)
    notes = []
    if language_status == "warn":
        notes.append("language differs from extracted source text")
    if content_status == "warn":
        notes.append("possible large omission or alteration")
    if not notes:
        notes.append("automatic heuristic check")

    return result_row(
        source_path=source_path,
        markdown_path=markdown_path,
        source_kind=source_kind,
        source_language=source_language,
        markdown_language=markdown_language,
        language_status=language_status,
        source_chars=source_chars,
        markdown_chars=markdown_chars,
        char_ratio=f"{char_ratio_value:.3f}",
        overlap=f"{overlap_value:.3f}" if overlap_value is not None else "",
        content_status=content_status,
        notes="; ".join(notes),
    )


def result_row(
    *,
    source_path: Path,
    markdown_path: Path,
    source_kind: str,
    source_language: str,
    markdown_language: str,
    language_status: str,
    source_chars: int,
    markdown_chars: int,
    char_ratio: str,
    overlap: str,
    content_status: str,
    notes: str,
) -> dict[str, str]:
    return {
        "path": str(source_path),
        "markdown_path": str(markdown_path),
        "source_kind": source_kind,
        "source_language": source_language,
        "markdown_language": markdown_language,
        "language_status": language_status,
        "source_chars": str(source_chars),
        "markdown_chars": str(markdown_chars),
        "char_ratio": char_ratio,
        "overlap": overlap,
        "content_status": content_status,
        "notes": notes,
    }


def inspect_pair_with_timeout(
    source_path: Path, markdown_path: Path, timeout_seconds: int
) -> dict[str, str]:
    queue: multiprocessing.Queue[dict[str, str]] = multiprocessing.Queue(maxsize=1)
    process = multiprocessing.Process(
        target=inspect_pair_worker,
        args=(str(source_path), str(markdown_path), queue),
    )
    process.start()
    process.join(timeout_seconds)
    if process.is_alive():
        process.terminate()
        process.join(5)
        if process.is_alive():
            process.kill()
            process.join()
        markdown = read_text(markdown_path) if markdown_path.exists() else ""
        return result_row(
            source_path=source_path,
            markdown_path=markdown_path,
            source_kind=source_path.suffix.lower().lstrip(".") or "unknown",
            source_language="unknown",
            markdown_language=classify_language(markdown),
            language_status="unknown",
            source_chars=0,
            markdown_chars=len(markdown),
            char_ratio="",
            overlap="",
            content_status="unknown",
            notes=f"source inspection exceeded {timeout_seconds} seconds",
        )
    try:
        return queue.get_nowait()
    except Empty:
        markdown = read_text(markdown_path) if markdown_path.exists() else ""
        return result_row(
            source_path=source_path,
            markdown_path=markdown_path,
            source_kind=source_path.suffix.lower().lstrip(".") or "unknown",
            source_language="unknown",
            markdown_language=classify_language(markdown),
            language_status="unknown",
            source_chars=0,
            markdown_chars=len(markdown),
            char_ratio="",
            overlap="",
            content_status="unknown",
            notes=f"inspection worker exited without result; exitcode={process.exitcode}",
        )


def inspect_pair_worker(
    source_path_value: str,
    markdown_path_value: str,
    queue: multiprocessing.Queue[dict[str, str]],
) -> None:
    queue.put(inspect_pair(Path(source_path_value), Path(markdown_path_value)))


def extract_source_text(path: Path) -> tuple[str, str]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        from pdfminer.high_level import extract_text

        return "pdfminer", extract_text(str(path)) or ""
    if suffix == ".docx":
        import mammoth

        with path.open("rb") as file:
            return "mammoth-docx", mammoth.extract_raw_text(file).value
    if suffix in {".pptx", ".potx"}:
        return "python-pptx", extract_pptx_text(path)
    if suffix == ".xlsx":
        return "openpyxl", extract_xlsx_text(path)
    if suffix == ".xls":
        return "xlrd", extract_xls_text(path)
    if suffix == ".zip":
        return "zip", extract_zip_text(path)
    if suffix in {".gz", ".tgz"} or path.name.endswith(".tar.gz"):
        return "tar", extract_tar_text(path)
    if suffix in TEXT_EXTENSIONS:
        return "text", read_text(path)
    try:
        text = read_text(path)
    except UnicodeDecodeError:
        return "unsupported", ""
    return "text-fallback", text


def extract_pptx_text(path: Path) -> str:
    from pptx import Presentation

    presentation = Presentation(str(path))
    texts: list[str] = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                texts.append(shape.text)
    return "\n".join(texts)


def extract_xlsx_text(path: Path) -> str:
    import openpyxl

    workbook = openpyxl.load_workbook(path, read_only=True, data_only=True)
    values: list[str] = []
    for sheet in workbook.worksheets:
        values.append(sheet.title)
        for row in sheet.iter_rows(values_only=True):
            values.extend(str(value) for value in row if value is not None)
    workbook.close()
    return "\n".join(values)


def extract_xls_text(path: Path) -> str:
    import xlrd

    workbook = xlrd.open_workbook(path)
    values: list[str] = []
    for sheet in workbook.sheets():
        values.append(sheet.name)
        for row_index in range(sheet.nrows):
            values.extend(str(value) for value in sheet.row_values(row_index) if value not in {"", None})
    return "\n".join(values)


def extract_zip_text(path: Path) -> str:
    texts: list[str] = []
    with zipfile.ZipFile(path) as archive:
        for name in archive.namelist()[:500]:
            texts.append(name)
            if is_text_member(name):
                with archive.open(name) as file:
                    texts.append(file.read(200_000).decode("utf-8", errors="replace"))
            elif name.endswith(".xml"):
                with archive.open(name) as file:
                    xml_text = file.read(200_000).decode("utf-8", errors="replace")
                texts.append(extract_xml_text(xml_text))
    return "\n".join(texts)


def extract_tar_text(path: Path) -> str:
    texts: list[str] = []
    with tarfile.open(path) as archive:
        for member in archive.getmembers()[:500]:
            texts.append(member.name)
            if member.isfile() and is_text_member(member.name):
                file = archive.extractfile(member)
                if file is not None:
                    texts.append(file.read(200_000).decode("utf-8", errors="replace"))
    return "\n".join(texts)


def extract_xml_text(xml_text: str) -> str:
    try:
        root = ElementTree.fromstring(xml_text)
    except ElementTree.ParseError:
        return xml_text
    return "\n".join(text.strip() for text in root.itertext() if text.strip())


def is_text_member(name: str) -> bool:
    return Path(name).suffix.lower() in TEXT_EXTENSIONS


def classify_language(text: str) -> str:
    sample = text[:200_000]
    if len(sample.strip()) < 80:
        return "unknown"
    japanese = len(re.findall(r"[\u3040-\u30ff\u3400-\u9fff]", sample))
    latin = len(re.findall(r"[A-Za-z]", sample))
    total = japanese + latin
    if total < 40:
        return "unknown"
    japanese_ratio = japanese / total
    latin_ratio = latin / total
    if japanese_ratio >= 0.25 and latin_ratio >= 0.25:
        return "mixed-ja-latin"
    if japanese_ratio >= 0.25:
        return "ja"
    if latin_ratio >= 0.70:
        return "latin"
    return "other"


def compare_language(source_language: str, markdown_language: str) -> str:
    if source_language == "unknown" or markdown_language == "unknown":
        return "unknown"
    if source_language == markdown_language:
        return "pass"
    if "mixed" in {source_language, markdown_language}:
        source_parts = set(source_language.split("-"))
        markdown_parts = set(markdown_language.split("-"))
        if source_parts & markdown_parts:
            return "pass"
    return "warn"


def compare_content(
    source_chars: int,
    markdown_chars: int,
    char_ratio: float,
    overlap: float | None,
) -> str:
    if source_chars < 200:
        return "unknown"
    if markdown_chars < 200:
        return "warn"
    if char_ratio < 0.45:
        return "warn"
    if overlap is not None and overlap < 0.35:
        return "warn"
    return "pass"


def overlap_score(source_text: str, markdown: str) -> float | None:
    source_norm = sample_for_overlap(normalize_for_overlap(source_text))
    markdown_norm = sample_for_overlap(normalize_for_overlap(markdown))
    if len(source_norm) < 120 or len(markdown_norm) < 120:
        return None
    source_shingles = shingles(source_norm, 12)
    markdown_shingles = shingles(markdown_norm, 12)
    if not source_shingles:
        return None
    return len(source_shingles & markdown_shingles) / len(source_shingles)


def normalize_for_overlap(text: str) -> str:
    normalized = unicodedata.normalize("NFKC", text).lower()
    return re.sub(r"\s+", "", normalized)


def sample_for_overlap(text: str, limit: int = 200_000) -> str:
    if len(text) <= limit:
        return text
    part = limit // 3
    middle_start = max(0, len(text) // 2 - part // 2)
    return text[:part] + text[middle_start : middle_start + part] + text[-part:]


def shingles(text: str, size: int) -> set[str]:
    return {hashlib.sha1(text[index : index + size].encode("utf-8")).hexdigest() for index in range(len(text) - size + 1)}


def write_tsv(out: Path, rows: list[dict[str, str]]) -> None:
    fieldnames = [
        "path",
        "markdown_path",
        "source_kind",
        "source_language",
        "markdown_language",
        "language_status",
        "source_chars",
        "markdown_chars",
        "char_ratio",
        "overlap",
        "content_status",
        "notes",
    ]
    with out.open("w", encoding="utf-8", newline="") as file:
        writer = csv.DictWriter(file, delimiter="\t", fieldnames=fieldnames)
        writer.writeheader()
        writer.writerows(rows)


def append_report(
    report: Path,
    inventory: Path,
    markdown_root: Path,
    inspection_tsv: Path,
    inventory_rows: list[dict[str, str]],
    inspected_rows: list[dict[str, str]],
) -> None:
    report.parent.mkdir(parents=True, exist_ok=True)
    now = dt.datetime.now(dt.UTC).astimezone().isoformat(timespec="seconds")
    status_counts = count_values(inventory_rows, "status")
    language_counts = count_values(inspected_rows, "language_status")
    content_counts = count_values(inspected_rows, "content_status")
    warnings = [
        row
        for row in inspected_rows
        if row["language_status"] == "warn" or row["content_status"] == "warn"
    ]
    unknowns = [
        row
        for row in inspected_rows
        if row["language_status"] == "unknown" or row["content_status"] == "unknown"
    ]

    section = [
        "",
        f"## {tool_name_from_path(inventory)} Output Inspection - {now}",
        "",
        "### Scope",
        "",
        f"- Inventory: `{inventory}`",
        f"- Markdown root: `{markdown_root}`",
        f"- Inspection TSV: `{inspection_tsv}`",
        f"- Inventory rows: {len(inventory_rows)}",
        f"- Generated Markdown inspected: {len(inspected_rows)}",
        f"- Inventory status counts: {format_counts(status_counts)}",
        "",
        "### Method",
        "",
        "- Extracted source text independently where practical: pdfminer for PDF, mammoth for DOCX, python-pptx for PPTX, openpyxl/xlrd for spreadsheets, and direct reads for text-like files.",
        "- Compared source and Markdown language with script-ratio heuristics: `ja`, `latin`, `mixed-ja-latin`, `other`, or `unknown`.",
        "- Checked large omissions or alteration risk with Markdown/source character ratio and normalized 12-character shingle overlap.",
        "- This is an automatic heuristic inspection. Rows marked `warn` or `unknown` need manual review before being treated as definitive quality judgments.",
        "",
        "### Summary",
        "",
        f"- Language status: {format_counts(language_counts)}",
        f"- Content status: {format_counts(content_counts)}",
        f"- Needs review: {len(warnings)} warn rows, {len(unknowns)} unknown rows",
        "",
        "### Warning Rows",
        "",
        "| Status | Source | Markdown | Ratio | Overlap | Notes |",
        "| ------ | ------ | -------- | ----- | ------- | ----- |",
    ]
    section.extend(format_issue_row(row) for row in warnings[:50])
    if len(warnings) > 50:
        section.append(f"| truncated | {len(warnings) - 50} more rows | | | | see TSV |")
    section.extend(
        [
            "",
            "### Unknown Rows",
            "",
            "| Source | Markdown | Source kind | Notes |",
            "| ------ | -------- | ----------- | ----- |",
        ]
    )
    section.extend(format_unknown_row(row) for row in unknowns[:50])
    if len(unknowns) > 50:
        section.append(f"| truncated | {len(unknowns) - 50} more rows | | see TSV |")
    section.append("")
    with report.open("a", encoding="utf-8") as file:
        file.write("\n".join(section))


def tool_name_from_path(path: Path) -> str:
    name = path.stem
    if name.endswith("-inventory"):
        name = name[: -len("-inventory")]
    return name or "conversion"


def count_values(rows: list[dict[str, str]], key: str) -> dict[str, int]:
    counts: dict[str, int] = {}
    for row in rows:
        counts[row[key]] = counts.get(row[key], 0) + 1
    return dict(sorted(counts.items()))


def format_counts(counts: dict[str, int]) -> str:
    return ", ".join(f"`{key}` {value}" for key, value in counts.items()) or "none"


def format_issue_row(row: dict[str, str]) -> str:
    status = f"{row['language_status']}/{row['content_status']}"
    return (
        f"| {escape_cell(status)} | {escape_cell(Path(row['path']).name)} | "
        f"{escape_cell(Path(row['markdown_path']).name)} | {row['char_ratio']} | "
        f"{row['overlap']} | {escape_cell(row['notes'])} |"
    )


def format_unknown_row(row: dict[str, str]) -> str:
    return (
        f"| {escape_cell(Path(row['path']).name)} | "
        f"{escape_cell(Path(row['markdown_path']).name)} | "
        f"{escape_cell(row['source_kind'])} | {escape_cell(row['notes'])} |"
    )


def escape_cell(value: str) -> str:
    return value.replace("|", "\\|").replace("\n", " ")


def read_text(path: Path) -> str:
    return path.read_text(encoding="utf-8", errors="replace")


def safe_name(value: str) -> str:
    stem = "".join(
        character if character.isascii() and (character.isalnum() or character in "-_") else "_"
        for character in value
    ).strip("_")
    digest = hashlib.sha256(value.encode("utf-8")).hexdigest()[:16]
    if not stem:
        stem = "file"
    return f"{stem[:120]}_{digest}"


if __name__ == "__main__":
    raise SystemExit(main())
